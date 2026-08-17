import re
from pathlib import Path
from typing import List, Dict, Any

class DocumentLoader:
    """
    Multi-format document parser that extracts clean text from academic files
    while preserving page-level and slide-level metadata.
    """

    @staticmethod
    def clean_academic_text(text: str) -> str:
        """Clean raw extracted text while preserving mathematical symbols and formatting."""
        if not text:
            return ""
        # Fix hyphenated line breaks (e.g., "algo-\nrithm" -> "algorithm")
        text = re.sub(r'(\w+)-\n(\w+)', r'\1\2', text)
        # Normalize multiple horizontal whitespaces
        text = re.sub(r'[ \t]+', ' ', text)
        # Normalize excessive newlines (max 2 consecutive newlines for paragraph separation)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    @classmethod
    def load_pdf(cls, file_path: Path) -> List[Dict[str, Any]]:
        """Load text page-by-page from a PDF file."""
        pages_data = []
        file_name = file_path.name
        
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(file_path))
            total_pages = len(reader.pages)
            
            for page_idx, page in enumerate(reader.pages, start=1):
                raw_text = page.extract_text() or ""
                cleaned = cls.clean_academic_text(raw_text)
                if cleaned:
                    pages_data.append({
                        "text": cleaned,
                        "source": file_name,
                        "page": page_idx,
                        "total_pages": total_pages,
                        "doc_type": "PDF",
                        "char_count": len(cleaned)
                    })
        except Exception as e:
            try:
                import pdfplumber
                with pdfplumber.open(str(file_path)) as pdf:
                    total_pages = len(pdf.pages)
                    for page_idx, page in enumerate(pdf.pages, start=1):
                        raw_text = page.extract_text() or ""
                        cleaned = cls.clean_academic_text(raw_text)
                        if cleaned:
                            pages_data.append({
                                "text": cleaned,
                                "source": file_name,
                                "page": page_idx,
                                "total_pages": total_pages,
                                "doc_type": "PDF",
                                "char_count": len(cleaned)
                            })
            except Exception as e2:
                print(f"[DocumentLoader] Error reading PDF {file_path}: {e} / {e2}")
                
        return pages_data

    @classmethod
    def load_docx(cls, file_path: Path) -> List[Dict[str, Any]]:
        """Load text from a DOCX file."""
        file_name = file_path.name
        pages_data = []
        try:
            import docx
            doc = docx.Document(str(file_path))
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())
            
            cleaned = cls.clean_academic_text("\n\n".join(full_text))
            if cleaned:
                pages_data.append({
                    "text": cleaned,
                    "source": file_name,
                    "page": 1,
                    "total_pages": 1,
                    "doc_type": "DOCX",
                    "char_count": len(cleaned)
                })
        except Exception as e:
            print(f"[DocumentLoader] Error reading DOCX {file_path}: {e}")
        return pages_data

    @classmethod
    def load_pptx(cls, file_path: Path) -> List[Dict[str, Any]]:
        """Load text slide-by-slide from a PPTX file."""
        file_name = file_path.name
        pages_data = []
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            total_slides = len(prs.slides)
            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                cleaned = cls.clean_academic_text("\n".join(slide_texts))
                if cleaned:
                    pages_data.append({
                        "text": cleaned,
                        "source": file_name,
                        "page": slide_idx,
                        "total_pages": total_slides,
                        "doc_type": "PPTX",
                        "char_count": len(cleaned)
                    })
        except Exception as e:
            print(f"[DocumentLoader] Error reading PPTX {file_path}: {e}")
        return pages_data

    @classmethod
    def load_txt(cls, file_path: Path) -> List[Dict[str, Any]]:
        """Load plain text or markdown file, splitting into simulated pages if large or delimited."""
        file_name = file_path.name
        pages_data = []
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            
            cleaned = cls.clean_academic_text(content)
            if not cleaned:
                return []
            
            # Check for explicit '--- Page \d+ ---' markers
            page_splits = re.split(r'---+\s*Page\s*\d+\s*---+', cleaned, flags=re.IGNORECASE)
            page_splits = [p.strip() for p in page_splits if p.strip()]
            
            if len(page_splits) > 1:
                total_pages = len(page_splits)
                for idx, p_text in enumerate(page_splits, start=1):
                    pages_data.append({
                        "text": p_text,
                        "source": file_name,
                        "page": idx,
                        "total_pages": total_pages,
                        "doc_type": "TXT",
                        "char_count": len(p_text)
                    })
            elif "\x0c" in cleaned:
                pages = cleaned.split("\x0c")
                total_pages = len(pages)
                for idx, p_text in enumerate(pages, start=1):
                    p_clean = p_text.strip()
                    if p_clean:
                        pages_data.append({
                            "text": p_clean,
                            "source": file_name,
                            "page": idx,
                            "total_pages": total_pages,
                            "doc_type": "TXT",
                            "char_count": len(p_clean)
                        })
            else:
                approx_page_len = 2000
                if len(cleaned) > approx_page_len:
                    chunks = [cleaned[i:i+approx_page_len] for i in range(0, len(cleaned), approx_page_len)]
                    total_pages = len(chunks)
                    for idx, ch in enumerate(chunks, start=1):
                        pages_data.append({
                            "text": ch.strip(),
                            "source": file_name,
                            "page": idx,
                            "total_pages": total_pages,
                            "doc_type": "TXT",
                            "char_count": len(ch.strip())
                        })
                else:
                    pages_data.append({
                        "text": cleaned,
                        "source": file_name,
                        "page": 1,
                        "total_pages": 1,
                        "doc_type": "TXT",
                        "char_count": len(cleaned)
                    })
        except Exception as e:
            print(f"[DocumentLoader] Error reading TXT {file_path}: {e}")
        return pages_data

    @classmethod
    def load_document(cls, file_path: Path) -> List[Dict[str, Any]]:
        """Unified document loader routing to the appropriate parser based on file extension."""
        suffix = file_path.suffix.lower()
        if suffix == ".pdf":
            return cls.load_pdf(file_path)
        elif suffix in [".docx", ".doc"]:
            return cls.load_docx(file_path)
        elif suffix in [".pptx", ".ppt"]:
            return cls.load_pptx(file_path)
        elif suffix in [".txt", ".md", ".markdown"]:
            return cls.load_txt(file_path)
        else:
            return cls.load_txt(file_path)
